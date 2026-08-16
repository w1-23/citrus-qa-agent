#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
md2pptx.py — Markdown → PPTX 转换器（Citrus QA Agent 柑橘主题）
用法: python md2pptx.py [input.md] [output.pptx]
语法:
  ---                    分页
  # 标题                 页面标题（首行）
  ## 小节标题
  - 要点（**加粗** 前缀为要点词；`代码` 橙色等宽）
    - 子要点
  > 引用/强调（左侧橙条卡片）
  | a | b |              表格（首行表头）
  <!-- type: cover|thanks -->     页面类型
  <!-- layout: two-col -->        双栏布局（左要点右示意图）
  <!-- visual: xxx -->            示意图（pillar/cards/browser/arch/blocks/
                                   statusbar/sticky/drawers/funnel/threshold/
                                   toolbox/gauge/team/pipeline/constitution/
                                   dials/timeline/quadrant）
示意图会"消费"页面的 - 要点行作为卡片/节点文本（视类型而定）。
"""
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ── 主题 ──────────────────────────────────────────────
ORANGE = RGBColor(0xF9, 0x73, 0x16)
ORANGE_DK = RGBColor(0xC2, 0x4E, 0x06)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
DARK = RGBColor(0x1F, 0x29, 0x37)
CREAM = RGBColor(0xFF, 0xF8, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_ORANGE = RGBColor(0xFD, 0xE8, 0xD6)
LIGHT_GREEN = RGBColor(0xE4, 0xF7, 0xEC)
RED = RGBColor(0xDC, 0x26, 0x26)
FONT = "微软雅黑"

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.55)
TITLE_H = Inches(1.02)
BODY_TOP = Inches(1.28)
BODY_H = Inches(5.7)


def _set_font(run, size, color, bold=False, font=FONT, italic=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", font)


def add_rich(p, text, size, color, bold=False):
    """段内解析 **粗体** 与 `代码`。"""
    pat = re.compile(r"(\*\*.+?\*\*|`[^`]+`)")
    pos = 0
    for m in pat.finditer(text):
        if m.start() > pos:
            r = p.add_run()
            r.text = text[pos:m.start()]
            _set_font(r, size, color, bold)
        seg = m.group(0)
        r = p.add_run()
        if seg.startswith("**"):
            r.text = seg[2:-2]
            _set_font(r, size, color, True)
        else:
            r.text = seg[1:-1]
            _set_font(r, size - 2, ORANGE, False, "Consolas")
        pos = m.end()
    if pos < len(text):
        r = p.add_run()
        r.text = text[pos:]
        _set_font(r, size, color, bold)


def add_box(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
            radius=0.06, line_w=1.0):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def add_text(slide, x, y, w, h, text, size=14, color=DARK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, line_spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    lines = str(text).split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing != 1.0:
            p.line_spacing = line_spacing
        add_rich(p, ln, size, color, bold)
    return tb


def add_arrow(slide, x1, y1, x2, y2, color=GREY, w=2.0):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(w)
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    return conn


def chip(slide, x, y, text, fill=ORANGE, color=WHITE, size=12, w=Inches(2.3), h=Inches(0.42)):
    add_box(slide, x, y, w, h, fill, radius=0.5)
    add_text(slide, x, y, w, h, text, size=size, color=color, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)


# ── 解析 ──────────────────────────────────────────────
def parse_pages(md: str):
    pages = []
    for block in md.split("\n---\n"):
        title, subs, dirs = "", [], {}
        for raw in block.strip("\n").splitlines():
            s = raw.strip()
            if not s:
                continue
            m = re.match(r"<!--\s*(\w+)(?::\s*([\w,.\- ]+))?\s*-->", s)
            if m:
                dirs[m.group(1)] = (m.group(2) or "").strip() or True
                continue
            if s.startswith("# ") and not title:
                title = s[2:].strip()
                continue
            if s.startswith("## "):
                subs.append(("h2", s[3:].strip()))
            elif s.startswith("> "):
                subs.append(("quote", s[2:].strip()))
            elif re.match(r"^\s*- ", raw):
                subs.append(("li", s[2:].strip()))
            elif re.match(r"^\s+[-*] ", raw):
                subs.append(("li2", re.sub(r"^\s+[-*] ", "", raw).strip()))
            elif s.startswith("|"):
                subs.append(("table", s))
            else:
                subs.append(("para", s))
        if title or subs:
            pages.append({"title": title, "subs": subs, "dir": dirs})
    return pages


def lis(page):
    return [c for k, c in page["subs"] if k == "li"]


def paras(page):
    return " ".join(c for k, c in page["subs"] if k == "para")


def split_li(li_text):
    """'- **标题**：描述' → (标题, 描述)"""
    m = re.match(r"\*\*(.+?)\*\*\s*[:：]?\s*(.*)", li_text)
    if m:
        return m.group(1), m.group(2)
    return li_text, ""


# ── 页面渲染 ──────────────────────────────────────────
def background(slide, color=CREAM):
    add_box(slide, 0, 0, W, H, color, shape=MSO_SHAPE.RECTANGLE)


def title_bar(slide, title, page_no, total):
    add_box(slide, 0, 0, W, TITLE_H, ORANGE, shape=MSO_SHAPE.RECTANGLE)
    add_box(slide, 0, TITLE_H - Inches(0.09), W, Inches(0.09), GREEN, shape=MSO_SHAPE.RECTANGLE)
    add_text(slide, MARGIN, 0, Inches(9.5), TITLE_H, title, size=26, color=WHITE,
             bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(10.3), 0, Inches(2.5), TITLE_H,
             f"Citrus QA Agent · {page_no}/{total}", size=11, color=WHITE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def render_items(slide, page, x, y, w, h, skip_li=False, size=15):
    cy = y
    for kind, content in page["subs"]:
        if kind in ("li", "li2") and skip_li:
            continue
        if kind == "h2":
            add_text(slide, x, cy, w, Inches(0.34), content, size=size - 1, color=ORANGE_DK, bold=True)
            cy += Inches(0.40)
        elif kind == "li":
            add_text(slide, x, cy, Inches(0.22), Inches(0.4), "●", size=size - 2, color=ORANGE, bold=True)
            add_text(slide, x + Inches(0.26), cy, w - Inches(0.26), Inches(0.4), content,
                     size=size, color=DARK)
            cy += Inches(0.44)
        elif kind == "li2":
            add_text(slide, x + Inches(0.5), cy, w - Inches(0.5), Inches(0.4), "– " + content,
                     size=size - 2, color=GREY)
            cy += Inches(0.40)
        elif kind == "para":
            add_text(slide, x, cy, w, Inches(0.5), content, size=size + 1, color=DARK)
            cy += Inches(0.56)
        elif kind == "quote":
            add_box(slide, x, cy, Inches(0.07), Inches(0.5), ORANGE, shape=MSO_SHAPE.RECTANGLE)
            add_box(slide, x + Inches(0.07), cy, w - Inches(0.07), Inches(0.5),
                    LIGHT_ORANGE, radius=0.12)
            add_text(slide, x + Inches(0.22), cy, w - Inches(0.35), Inches(0.5), content,
                     size=size - 1, color=ORANGE_DK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            cy += Inches(0.62)
        elif kind == "table":
            cy = render_table(slide, content, x, cy, w)
    return cy


def render_table(slide, md_rows, x, y, w):
    rows = []
    for line in md_rows.splitlines():
        if line.strip().startswith("|"):
            cells = [c.strip() for c in line.strip().strip("|").split("|")]
            rows.append(cells)
    if not rows:
        return y
    n_r, n_c = len(rows), max(len(r) for r in rows)
    col_w = Emu(int(w / n_c))
    row_h = Inches(0.34)
    gt = slide.shapes.add_table(n_r, n_c, x, y, w, row_h * n_r).table
    for ci in range(n_c):
        gt.columns[ci].width = col_w
    for ri in range(n_r):
        gt.rows[ri].height = row_h
        for ci in range(n_c):
            cell = gt.cell(ri, ci)
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = ORANGE if ri == 0 else (LIGHT_ORANGE if ri % 2 else WHITE)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            txt = rows[ri][ci] if ci < len(rows[ri]) else ""
            add_rich(p, txt, 11, WHITE if ri == 0 else DARK, ri == 0)
    return y + row_h * n_r + Inches(0.14)


# ── 示意图 ────────────────────────────────────────────
VISUAL_CONSUMES_LI = {"pillar", "cards", "timeline", "quadrant"}


def draw_visual(slide, page, x, y, w, h):
    v = page["dir"].get("visual", "")
    name, _, arg = v.partition(":")
    fn = {
        "pillar": v_pillar, "cards": v_cards, "browser": v_browser, "arch": v_arch,
        "blocks": v_blocks, "statusbar": v_statusbar, "sticky": v_sticky,
        "drawers": v_drawers, "funnel": v_funnel, "threshold": v_threshold,
        "toolbox": v_toolbox, "gauge": v_gauge, "team": v_team, "pipeline": v_pipeline,
        "constitution": v_constitution, "dials": v_dials, "timeline": v_timeline,
        "quadrant": v_quadrant,
    }.get(name)
    if fn:
        fn(slide, page, x, y, w, h, arg)
    else:
        add_box(slide, x, y, w, h, LIGHT_ORANGE, radius=0.08)
        add_text(slide, x, y, w, h, f"示意图：{v or '(无)'}", size=14, color=ORANGE_DK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _card(slide, x, y, w, h, title, desc, fill=WHITE, tcolor=ORANGE_DK, title_size=13, desc_size=11):
    add_box(slide, x, y, w, h, fill, line=RGBColor(0xE8, 0xDD, 0xD0))
    add_text(slide, x + Inches(0.12), y + Inches(0.10), w - Inches(0.24), Inches(0.3),
             title, size=title_size, color=tcolor, bold=True)
    if desc:
        add_text(slide, x + Inches(0.12), y + Inches(0.42), w - Inches(0.24), h - Inches(0.5),
                 desc, size=desc_size, color=DARK)


def v_pillar(slide, page, x, y, w, h, arg):
    items = lis(page)[:3]
    labels = ["文献太多", "搜索引擎不够用", "通用 AI 会幻觉"]
    heights = [Inches(2.3), Inches(1.7), Inches(1.1)]
    colors = [GREEN, ORANGE, RED]
    n = 3
    pw = Emu(int(w / n) - Inches(0.25))
    for i in range(n):
        px = x + Emu(int(w / n) * i) + Inches(0.12)
        ph = heights[i]
        add_box(slide, px, y + h - ph - Inches(0.5), pw, ph, colors[i], radius=0.1)
        t, d = split_li(items[i]) if i < len(items) else (labels[i], "")
        add_text(slide, px, y + h - ph - Inches(0.78), pw, Inches(0.4), labels[i],
                 size=13, color=DARK, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, px, y + Inches(0.05), pw, Inches(0.5),
                 "①②③"[i] + ("  " + t if t else ""), size=12, color=GREY, align=PP_ALIGN.CENTER)


def v_cards(slide, page, x, y, w, h, arg):
    try:
        n = int(arg)
    except Exception:
        n = 3
    items = lis(page)
    cols = 2 if n > 3 else n
    rows_n = (n + cols - 1) // cols
    cw = Emu(int(w / cols) - Inches(0.22))
    ch = Emu(int((h - Inches(0.2)) / rows_n) - Inches(0.16))
    for i in range(n):
        cx = x + Emu(int(w / cols) * (i % cols)) + Inches(0.1)
        cy = y + Emu(int((h - Inches(0.2)) / rows_n) * (i // cols)) + Inches(0.08)
        t, d = split_li(items[i]) if i < len(items) else ("", "")
        _card(slide, cx, cy, cw, ch, t, d)


def v_browser(slide, page, x, y, w, h, arg):
    bw, bh = Emu(int(w * 0.96)), Inches(4.6)
    add_box(slide, x, y, bw, bh, WHITE, line=RGBColor(0xE0, 0xD5, 0xC8), radius=0.05)
    # 顶栏
    add_box(slide, x, y, bw, Inches(0.42), RGBColor(0xF1, 0xE8, 0xDD), radius=0.05)
    for i, c in enumerate([RED, ORANGE, GREEN]):
        add_box(slide, x + Inches(0.14 + i * 0.2), y + Inches(0.13), Inches(0.16), Inches(0.16), c, shape=MSO_SHAPE.OVAL)
    add_box(slide, x + Inches(0.85), y + Inches(0.09), bw - Inches(1.7), Inches(0.24),
            WHITE, radius=0.5)
    # 聊天气泡
    add_box(slide, x + Inches(0.3), y + Inches(0.7), Inches(4.0), Inches(0.62), LIGHT_ORANGE, radius=0.3)
    add_text(slide, x + Inches(0.45), y + Inches(0.76), Inches(3.8), Inches(0.5),
             "血橙为什么是红的？", size=11, color=DARK)
    add_box(slide, x + Inches(0.3), y + Inches(1.5), Inches(5.1), Inches(1.5), WHITE,
            line=RGBColor(0xE8, 0xDD, 0xD0), radius=0.12)
    add_text(slide, x + Inches(0.45), y + Inches(1.62), Inches(4.9), Inches(1.3),
             "红色主要来自**花青苷**在果皮与果肉的积累 [1][2]……\n\n证据：花色苷合成受 Ruby 基因调控 [3]", size=11, color=DARK)
    # 引用侧栏
    sx = x + Inches(5.75)
    add_text(slide, sx, y + Inches(0.6), Inches(1.7), Inches(0.3), "文献引用", size=11, color=GREY, bold=True)
    for i, t in enumerate(["花青苷积累机制 (2023)", "Ruby 基因调控 (2021)"]):
        add_box(slide, sx, y + Inches(0.95 + i * 0.62), Inches(1.75), Inches(0.52), LIGHT_GREEN, radius=0.15)
        add_text(slide, sx + Inches(0.1), y + Inches(1.02 + i * 0.62), Inches(1.55), Inches(0.4),
                 f"[{i+1}] {t}", size=9, color=DARK)


def v_arch(slide, page, x, y, w, h, arg):
    items = lis(page)[:3]
    names = ["LLM 大脑", "上下文 眼睛", "工具 手脚"]
    cw, ch = Inches(1.9), Inches(0.8)
    for i in range(3):
        cx = x + Emu(int(w / 3) * i) + Inches(0.25)
        t = split_li(items[i])[0] if i < len(items) else names[i]
        add_box(slide, cx, y, cw, ch, [ORANGE, GREEN, ORANGE_DK][i], radius=0.2)
        add_text(slide, cx, y, cw, ch, t, size=13, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < 2:
            add_arrow(slide, cx + cw + Inches(0.08), y + ch / 2, cx + cw + Inches(0.5), y + ch / 2, ORANGE)
    # 架构图：用户 → 主管 → 专员
    ay = y + Inches(1.5)
    add_text(slide, x, ay, w, Inches(0.3), "整体架构：主管 + 三个专员", size=13, color=ORANGE_DK, bold=True,
             align=PP_ALIGN.CENTER)
    ux, uy = x + Inches(0.1), ay + Inches(0.45)
    add_box(slide, ux, uy, Inches(1.5), Inches(0.7), GREY, radius=0.25)
    add_text(slide, ux, uy, Inches(1.5), Inches(0.7), "用户", size=13, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_arrow(slide, ux + Inches(1.5), uy + Inches(0.3), ux + Inches(2.3), uy + Inches(0.3), ORANGE, 2.5)
    mx, my = ux + Inches(2.35), uy
    add_box(slide, mx, my, Inches(2.2), Inches(0.7), ORANGE, radius=0.25)
    add_text(slide, mx, my, Inches(2.2), Inches(0.7), "Supervisor 主管", size=13, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    subs = ["🔍 检索专员", "✍️ 写作专员", "📊 分析专员"]
    sw = Inches(1.85)
    for i in range(3):
        sx_ = x + Inches(0.35) + Emu(int((w - Inches(0.7)) / 3) * i)
        add_arrow(slide, mx + Inches(0.9), my + Inches(0.7), sx_ + sw / 2, my + Inches(1.35), GREY, 2.0)
        add_box(slide, sx_, my + Inches(1.4), sw, Inches(0.66), WHITE, line=RGBColor(0xE8, 0xDD, 0xD0), radius=0.25)
        add_text(slide, sx_, my + Inches(1.4), sw, Inches(0.66), subs[i], size=12, color=DARK, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x, my + Inches(2.3), w, Inches(0.9),
             "RAG 知识库 · 记忆 · 工具包", size=12, color=GREY, align=PP_ALIGN.CENTER)


def v_blocks(slide, page, x, y, w, h, arg):
    n_static, n_dyn = 3, 2
    bw, bh = Inches(1.15), Inches(1.5)
    gap = Inches(0.14)
    total = n_static * (bw + gap) + n_dyn * (bw + gap)
    sx = x + (w - total) / 2
    for i in range(n_static):
        add_box(slide, sx + Emu(int((bw + gap) * i)), y, bw, bh, ORANGE, radius=0.1)
        add_text(slide, sx + Emu(int((bw + gap) * i)), y + Inches(0.55), bw, Inches(0.5), "🔒", size=20,
                 align=PP_ALIGN.CENTER)
    for i in range(n_dyn):
        add_box(slide, sx + Emu(int((bw + gap) * (i + n_static))), y, bw, bh, RGBColor(0xE8, 0xDD, 0xD0), radius=0.1)
        add_text(slide, sx + Emu(int((bw + gap) * (i + n_static))), y + Inches(0.55), bw, Inches(0.5),
                 "+", size=24, color=GREY, align=PP_ALIGN.CENTER)
    add_text(slide, sx, y + bh + Inches(0.12), Emu(int(total * 0.55)), Inches(0.4),
             "静态前缀 · 字节级稳定", size=12, color=ORANGE_DK, bold=True)
    add_text(slide, sx + Emu(int(total * 0.55)), y + bh + Inches(0.12), Emu(int(total * 0.45)), Inches(0.4),
             "动态内容 · 追加尾部", size=12, color=GREY)
    chip(slide, sx + Inches(0.2), y + bh + Inches(0.62), "命中 DeepSeek 上下文缓存", GREEN, size=11, w=Inches(3.4))


def v_statusbar(slide, page, x, y, w, h, arg):
    # 压缩示意
    add_text(slide, x, y, w, Inches(0.32), "分级压缩", size=13, color=ORANGE_DK, bold=True)
    bw = Emu(int(w * 0.78))
    add_box(slide, x, y + Inches(0.4), bw, Inches(0.6), RGBColor(0xE8, 0xDD, 0xD0), radius=0.1)
    add_text(slide, x + Inches(0.15), y + Inches(0.42), Inches(2.5), Inches(0.5), "100% 未整理", size=11, color=GREY)
    add_arrow(slide, x + bw + Inches(0.05), y + Inches(0.7), x + bw + Inches(0.55), y + Inches(0.7), ORANGE, 2.5)
    add_box(slide, x + bw + Inches(0.65), y + Inches(0.4), Emu(int(bw * 0.5)), Inches(0.6), ORANGE, radius=0.1)
    add_text(slide, x + bw + Inches(0.8), y + Inches(0.42), Inches(2.5), Inches(0.5), "→ 50% 压缩后", size=11, color=WHITE, bold=True)
    # 保护名单
    add_text(slide, x, y + Inches(1.25), w, Inches(0.3), "保护名单（绝不压缩）", size=12, color=GREY, bold=True)
    for i, t in enumerate(["最近 3 轮问答", "被引证据全文", "DOI 编号"]):
        chip(slide, x + Emu(int((w - Inches(0.6)) / 3) * i), y + Inches(1.6), t, LIGHT_GREEN, DARK, size=11,
             w=Emu(int((w - Inches(0.6)) / 3) - Inches(0.1)))
    # 状态栏
    add_text(slide, x, y + Inches(2.3), w, Inches(0.3), "<agent_status> 状态栏（代码注入）", size=12, color=GREY, bold=True)
    add_box(slide, x, y + Inches(2.62), Emu(int(w * 0.95)), Inches(0.62), DARK, radius=0.2)
    add_text(slide, x + Inches(0.2), y + Inches(2.62), Emu(int(w * 0.9)), Inches(0.62),
             "<agent_status> 轮次:2/4 · 已检索文献:6 · 预算剩余:68%", size=12, color=RGBColor(0x7E, 0xE0, 0x9B),
             bold=True, anchor=MSO_ANCHOR.MIDDLE)


def v_sticky(slide, page, x, y, w, h, arg):
    notes = ["用户偏好中文", "研究血橙花色苷", "置信 ≥0.8 才上墙", "≤8 张 · 自动淘汰"]
    colors = [RGBColor(0xFF, 0xE0, 0xB2), RGBColor(0xFF, 0xF3, 0xCD), RGBColor(0xE4, 0xF7, 0xEC),
              RGBColor(0xFD, 0xE8, 0xD6)]
    nw, nh = Inches(1.9), Inches(1.15)
    for i in range(4):
        nx = x + Inches(0.15) + Emu(int((w - Inches(0.5)) / 2) * (i % 2))
        ny = y + Inches(0.2) + Inches(1.35) * (i // 2)
        sp = add_box(slide, nx, ny, nw, nh, colors[i], radius=0.08)
        sp.rotation = 6 if i % 2 == 0 else -5
        add_text(slide, nx + Inches(0.15), ny + Inches(0.3), nw - Inches(0.3), nh - Inches(0.4),
                 notes[i], size=12, color=DARK, bold=True)
    # 日志本
    add_box(slide, x + Inches(0.15), y + Inches(2.9), Inches(3.4), Inches(1.5), WHITE,
            line=RGBColor(0xE8, 0xDD, 0xD0), radius=0.06)
    add_text(slide, x + Inches(0.3), y + Inches(3.05), Inches(3.1), Inches(0.3),
             "📔 长期事实库 LTM", size=13, color=ORANGE_DK, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(3.5), Inches(3.1), Inches(0.9),
             "ADD-only 只追加不覆盖\n置信 <0.5 拒绝写入\n语义 × 时间衰减召回", size=11, color=GREY)


def v_drawers(slide, page, x, y, w, h, arg):
    add_text(slide, x, y, w, Inches(0.32), "记忆按会话分抽屉", size=13, color=ORANGE_DK, bold=True)
    labels = ["会话 A 的记忆", "会话 B 的记忆", "全局偏好（跨会话）"]
    fills = [RGBColor(0xE8, 0xDD, 0xD0), RGBColor(0xE8, 0xDD, 0xD0), ORANGE]
    for i in range(3):
        add_box(slide, x + Inches(0.6), y + Inches(0.45 + i * 1.15), Inches(3.6), Inches(0.95),
                fills[i], radius=0.1)
        add_text(slide, x + Inches(0.85), y + Inches(0.45 + i * 1.15), Inches(3.2), Inches(0.95),
                 labels[i], size=13, color=WHITE if i == 2 else DARK, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(4.6), y + Inches(1.1), Inches(1.9), Inches(2.5),
             "「综述用中文」→ 全局\n新会话也生效\n\n会话细节 → 不跨域乱串\n\n置信 ≥0.9 才全局共享",
             size=11, color=GREY)


def v_funnel(slide, page, x, y, w, h, arg):
    steps = [("PDF / txt / md 文献", "导入"), ("分块 Chunking", "11.9 万条片段"),
             ("向量化嵌入", "e5 多语言模型"), ("LanceDB 向量库", "百万级 · 热更新")]
    cw0 = Emu(int(w * 0.85))
    for i, (t, d) in enumerate(steps):
        cw = Emu(int(cw0 * (1 - i * 0.17)))
        cx = x + (w - cw) / 2
        cy = y + Inches(0.15 + i * 1.15)
        add_box(slide, cx, cy, cw, Inches(0.9), [ORANGE, GREEN, ORANGE_DK, RGBColor(0xE8, 0xDD, 0xD0)][i], radius=0.12)
        add_text(slide, cx, cy + Inches(0.12), cw, Inches(0.4), t, size=13, color=WHITE if i < 3 else DARK,
                 bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, cx, cy + Inches(0.5), cw, Inches(0.3), d, size=10, color=WHITE if i < 3 else GREY,
                 align=PP_ALIGN.CENTER)
        if i < 3:
            add_arrow(slide, cx + cw / 2, cy + Inches(0.9), cx + cw / 2, cy + Inches(1.08), GREY)


def v_threshold(slide, page, x, y, w, h, arg):
    # 直方图：通过(绿/橙) vs 过滤(灰)，阈值线
    vals = [0.30, 0.42, 0.55, 0.68, 0.82, 0.90, 0.72, 0.58, 0.38, 0.26]
    thresh = 0.55
    n = len(vals)
    bw = Emu(int(w / (n + 1.5)))
    base = y + Inches(2.6)
    for i, v in enumerate(vals):
        bh = Emu(int(Inches(2.2) * v))
        passed = v >= thresh
        add_box(slide, x + Emu(int(bw * (i + 0.6))), base - bh, Emu(int(bw * 0.6)), bh,
                GREEN if passed else RGBColor(0xD8, 0xD2, 0xC8), radius=0.08)
    # 阈值线
    line_y = base - Emu(int(Inches(2.2) * thresh))
    add_box(slide, x, line_y, w, Inches(0.03), RED, shape=MSO_SHAPE.RECTANGLE)
    add_text(slide, x + Inches(0.1), line_y - Inches(0.3), Inches(3.2), Inches(0.3),
             f"动态阈值 = max(0.25 地板, top×0.6)", size=11, color=RED, bold=True)
    add_text(slide, x, base + Inches(0.12), w, Inches(0.6),
             "分数分布写入检索日志 → 按统计分布更新地板（自适应）", size=11, color=GREY)
    add_text(slide, x, y, w, Inches(0.32), "两阶段检索：海选（BM25+向量）→ 面试（重排精排）",
             size=13, color=ORANGE_DK, bold=True)


def v_toolbox(slide, page, x, y, w, h, arg):
    tools = [("🔍 检索", "本地 RAG + 学术库"), ("📖 读文件", "PDF / MD / CSV"),
             ("✍️ 写文件", "沙箱 workspace/output"), ("📊 分析", "统计检验")]
    for i, (t, d) in enumerate(tools):
        cx = x + Inches(0.1) + Emu(int((w - Inches(0.4)) / 2) * (i % 2))
        cy = y + Inches(0.1) + Inches(1.5) * (i // 2)
        _card(slide, cx, cy, Emu(int((w - Inches(0.4)) / 2) - Inches(0.08)), Inches(1.3), t, d,
              title_size=15, desc_size=11)
    # 审批卡片
    px, py = x + Inches(2.9), y + Inches(3.0)
    add_box(slide, px, py, Inches(2.9), Inches(1.7), WHITE, line=ORANGE, line_w=2.0, radius=0.1)
    add_text(slide, px + Inches(0.15), py + Inches(0.12), Inches(2.6), Inches(0.3),
             "⚠️ 权限审批", size=12, color=RED, bold=True)
    add_text(slide, px + Inches(0.15), py + Inches(0.45), Inches(2.6), Inches(0.6),
             "请求写入文件 output/report.md", size=11, color=DARK)
    chip(slide, px + Inches(0.15), py + Inches(1.15), "允许本次", ORANGE, size=10, w=Inches(1.2), h=Inches(0.38))
    chip(slide, px + Inches(1.55), py + Inches(1.15), "拒绝", GREY, size=10, w=Inches(1.2), h=Inches(0.38))


def v_gauge(slide, page, x, y, w, h, arg):
    meters = [("检索预算", "6/6 用尽", ORANGE), ("边际收敛", "≥6 篇 · 新增<25%", GREEN),
              ("熔断", "连续 3 次失败", RED), ("停止按钮", "一键取消", GREY)]
    for i, (t, d, c) in enumerate(meters):
        cx = x + Inches(0.1) + Emu(int((w - Inches(0.4)) / 2) * (i % 2))
        cy = y + Inches(0.1) + Inches(2.15) * (i // 2)
        add_box(slide, cx, cy, Inches(1.7), Inches(1.7), LIGHT_ORANGE, shape=MSO_SHAPE.OVAL)
        add_box(slide, cx + Inches(0.22), cy + Inches(0.22), Inches(1.26), Inches(1.26),
                CREAM, shape=MSO_SHAPE.OVAL)
        add_text(slide, cx, cy + Inches(0.6), Inches(1.7), Inches(0.5), t, size=13, color=DARK, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(slide, cx, cy + Inches(1.05), Inches(1.7), Inches(0.4), d, size=9, color=c,
                 align=PP_ALIGN.CENTER, bold=True)


def v_team(slide, page, x, y, w, h, arg):
    mx = x + w / 2 - Inches(1.1)
    add_box(slide, mx, y, Inches(2.2), Inches(0.75), ORANGE, radius=0.25)
    add_text(slide, mx, y, Inches(2.2), Inches(0.75), "Supervisor 主管", size=14, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    subs = ["🔍 检索专员", "✍️ 写作专员", "📊 分析专员"]
    sw = Inches(1.95)
    for i in range(3):
        sx_ = x + Inches(0.3) + Emu(int((w - Inches(0.6)) / 3) * i)
        add_arrow(slide, mx + Inches(1.1), y + Inches(0.75), sx_ + sw / 2, y + Inches(1.45), ORANGE, 2.5)
        add_box(slide, sx_, y + Inches(1.5), sw, Inches(0.7), WHITE, line=RGBColor(0xE8, 0xDD, 0xD0), radius=0.25)
        add_text(slide, sx_, y + Inches(1.5), sw, Inches(0.7), subs[i], size=13, color=DARK, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_box(slide, x + Inches(0.2), y + Inches(2.6), w - Inches(0.4), Inches(1.5), LIGHT_GREEN, radius=0.1)
    add_text(slide, x + Inches(0.4), y + Inches(2.75), w - Inches(0.8), Inches(1.3),
             "**与编码 Agent 的本质差异**：编码类 = 隔离 + 压缩传摘要；\n"
             "我们 = **证据全文共享**（结论 / 数值 / 机制细节只在原文）\n\n"
             "**历史证据账本**：追问时全额注入上一轮证据 → 免重新检索",
             size=12, color=DARK)


def v_pipeline(slide, page, x, y, w, h, arg):
    chain = paras(page)
    nodes = [s.strip() for s in chain.split("→") if s.strip()]
    if len(nodes) < 2:
        nodes = ["任务分类", "材料打包", "生成大纲", "大纲校验", "并行写章", "引用统一", "完整性校验", "原子发布"]
    n = len(nodes)
    nw, nh = Inches(1.28), Inches(1.15)
    gap = Inches(0.42)
    total = n * nw + (n - 1) * gap
    sx = x + (w - total) / 2
    add_text(slide, x, y - Inches(0.02), w, Inches(0.32), "写作流水线 Plan-Execute", size=13,
             color=ORANGE_DK, bold=True, align=PP_ALIGN.CENTER)
    for i, t in enumerate(nodes):
        nx = sx + Emu(int((nw + gap) * i))
        add_box(slide, nx, y + Inches(0.42), nw, nh, ORANGE if i in (2, 4, 7) else GREEN, radius=0.15)
        add_text(slide, nx, y + Inches(0.42), nw, Inches(0.3), f"{i+1}", size=10, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(slide, nx, y + Inches(0.72), nw, Inches(0.8), t, size=10, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER)
        if i < n - 1:
            add_arrow(slide, nx + nw, y + Inches(0.95), nx + nw + gap, y + Inches(0.95), GREY)
    add_text(slide, x, y + Inches(1.95), w, Inches(0.9),
             "原子发布：先写 .draft.md → 校验通过 → 一次性改名落盘 · 中断不留半成品 · 断点续传",
             size=11, color=GREY, align=PP_ALIGN.CENTER)


def v_constitution(slide, page, x, y, w, h, arg):
    add_box(slide, x, y, Inches(3.4), h, ORANGE, radius=0.06)
    add_text(slide, x, y + Inches(0.25), Inches(3.4), Inches(0.5), "232 项回归测试", size=16, color=WHITE,
             bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(0.85), Inches(3.4), Inches(0.4), "任何改动不许变红", size=12, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_box(slide, x + Inches(0.45), y + Inches(1.6), Inches(2.5), Inches(0.9), LIGHT_ORANGE, radius=0.15)
    add_text(slide, x + Inches(0.6), y + Inches(1.75), Inches(2.2), Inches(0.7), "✓ 232 passed", size=20,
             color=ORANGE_DK, bold=True, align=PP_ALIGN.CENTER)
    invs = ["INV-01 协议配对", "INV-02 生命周期有界", "INV-03 检索降级", "INV-04 上下文传播",
            "INV-05 输出路由", "INV-06 可观测回归", "INV-07 状态显式化", "INV-08 熔断与输入隔离",
            "INV-09 证据账本", "INV-10 存储全量·发送裁剪"]
    add_box(slide, x + Inches(3.7), y, w - Inches(3.7), h, WHITE, line=RGBColor(0xE8, 0xDD, 0xD0), radius=0.06)
    for i, t in enumerate(invs):
        add_text(slide, x + Inches(3.95), y + Inches(0.15 + i * 0.42), w - Inches(4.1), Inches(0.4),
                 ("● " if i < 4 else "○ ") + t, size=11, color=ORANGE_DK if i < 4 else GREY,
                 bold=(i < 4))
    add_text(slide, x + Inches(3.95), y + h - Inches(0.45), w - Inches(4.1), Inches(0.4),
             "每个不变量背后 = 一个真实事故（F1~F9 档案）", size=10, color=GREY, bold=True)


def v_dials(slide, page, x, y, w, h, arg):
    meters = [("缓存命中率", "82%", GREEN), ("检索平均耗时", "320ms", ORANGE),
              ("单请求工具调用", "≤6 次", ORANGE_DK)]
    for i, (t, v, c) in enumerate(meters):
        cx = x + Inches(0.1) + Emu(int((w - Inches(0.4)) / 3) * i)
        add_box(slide, cx, y + Inches(0.15), Inches(1.9), Inches(1.9), LIGHT_ORANGE, shape=MSO_SHAPE.OVAL)
        add_box(slide, cx + Inches(0.25), y + Inches(0.4), Inches(1.4), Inches(1.4),
                CREAM, shape=MSO_SHAPE.OVAL)
        add_text(slide, cx, y + Inches(0.75), Inches(1.9), Inches(0.5), v, size=18, color=c, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(slide, cx, y + Inches(1.35), Inches(1.9), Inches(0.4), t, size=11, color=DARK,
                 align=PP_ALIGN.CENTER)
    # 反馈循环
    add_text(slide, x, y + Inches(2.35), w, Inches(0.4), "反馈闭环：👍/👎 + 评论 → 本地库 → 改进",
             size=13, color=ORANGE_DK, bold=True, align=PP_ALIGN.CENTER)
    loop = ["👍/👎 反馈", "内容哈希幂等", "三级日志", "假完成检测"]
    for i, t in enumerate(loop):
        chip(slide, x + Emu(int((w - Inches(0.8)) / 4) * i), y + Inches(2.85), t, WHITE, DARK, size=11,
             w=Emu(int((w - Inches(0.8)) / 4) - Inches(0.12)), h=Inches(0.5))


def v_timeline(slide, page, x, y, w, h, arg):
    items = lis(page)[:6]
    line_x = x + Inches(0.55)
    add_box(slide, line_x, y + Inches(0.1), Inches(0.045), Inches(4.6), ORANGE, shape=MSO_SHAPE.RECTANGLE)
    labels = ["接单", "派活", "找材料", "写作", "交稿", "评价"]
    for i, it in enumerate(items):
        ny = y + Inches(0.15 + i * 0.78)
        add_box(slide, line_x - Inches(0.13), ny + Inches(0.16), Inches(0.3), Inches(0.3), GREEN, shape=MSO_SHAPE.OVAL)
        add_text(slide, line_x - Inches(0.07), ny + Inches(0.2), Inches(0.2), Inches(0.2), str(i + 1),
                 size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        t, d = split_li(it)
        add_box(slide, line_x + Inches(0.25), ny, Emu(int(w - Inches(0.8))), Inches(0.68), WHITE,
                line=RGBColor(0xE8, 0xDD, 0xD0), radius=0.12)
        add_text(slide, line_x + Inches(0.42), ny + Inches(0.07), Inches(1.4), Inches(0.3),
                 f"{labels[i]} · {t}", size=11, color=ORANGE_DK, bold=True)
        add_text(slide, line_x + Inches(1.9), ny + Inches(0.09), Emu(int(w - Inches(2.6))), Inches(0.5),
                 d, size=10, color=GREY)


def v_quadrant(slide, page, x, y, w, h, arg):
    items = lis(page)[:4]
    titles = ["Key 不出门", "数据全本地", "防提示注入", "日志脱敏"]
    for i in range(4):
        cx = x + Emu(int(w / 2) * (i % 2)) + Inches(0.08)
        cy = y + Emu(int(h / 2) * (i // 2)) + Inches(0.08)
        cw = Emu(int(w / 2) - Inches(0.24))
        ch = Emu(int(h / 2) - Inches(0.24))
        t, d = split_li(items[i]) if i < len(items) else (titles[i], "")
        add_box(slide, cx, cy, cw, ch, WHITE, line=RGBColor(0xE8, 0xDD, 0xD0), radius=0.08)
        add_text(slide, cx + Inches(0.2), cy + Inches(0.18), cw - Inches(0.4), Inches(0.4),
                 t, size=14, color=ORANGE_DK, bold=True)
        add_text(slide, cx + Inches(0.2), cy + Inches(0.65), cw - Inches(0.4), ch - Inches(0.8),
                 d, size=11, color=GREY)


# ── 组装 ──────────────────────────────────────────────
def render_cover(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide, ORANGE)
    add_box(slide, 0, H - Inches(0.18), W, Inches(0.18), GREEN, shape=MSO_SHAPE.RECTANGLE)
    add_box(slide, Inches(10.9), Inches(-1.4), Inches(4.4), Inches(4.4), ORANGE_DK, shape=MSO_SHAPE.OVAL)
    add_box(slide, Inches(11.6), Inches(-0.7), Inches(3.0), Inches(3.0), ORANGE, shape=MSO_SHAPE.OVAL)
    add_box(slide, Inches(12.25), Inches(-0.05), Inches(1.7), Inches(1.7), GREEN, shape=MSO_SHAPE.OVAL)
    add_text(slide, Inches(0.9), Inches(1.6), Inches(9.5), Inches(1.1), page["title"], size=52,
             color=WHITE, bold=True)
    add_text(slide, Inches(0.95), Inches(2.75), Inches(10), Inches(0.6),
             "你的柑橘科研问答助手", size=22, color=WHITE)
    chips = ["本地运行", "证据可溯", "自动写综述"]
    for i, c in enumerate(chips):
        chip(slide, Inches(0.95 + i * 2.5), Inches(3.7), c, GREEN, size=14, w=Inches(2.2), h=Inches(0.55))
    sub = page["dir"].get("subtitle", "")
    add_text(slide, Inches(0.95), Inches(5.0), Inches(10), Inches(0.5), str(sub), size=14, color=WHITE)
    add_text(slide, Inches(0.95), Inches(6.5), Inches(10), Inches(0.4),
             "AI-Agents-in-Depth 理论 · 完整工程落地", size=12, color=RGBColor(0xFF, 0xE0, 0xC0))


def render_thanks(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide)
    title_bar(slide, page["title"], total, total)
    items = [(k, c) for k, c in page["subs"] if k in ("li", "para", "quote")]
    cy = BODY_TOP + Inches(0.2)
    for k, it in items:
        t, d = split_li(it)
        if t and d:
            add_text(slide, MARGIN, cy, Inches(12.2), Inches(0.4), "▪ " + t, size=19, color=ORANGE_DK, bold=True)
            cy += Inches(0.5)
            add_text(slide, MARGIN + Inches(0.3), cy, Inches(11.8), Inches(0.4), d, size=15, color=DARK)
            cy += Inches(0.55)
        else:
            add_text(slide, MARGIN, cy, Inches(12.2), Inches(0.5), it, size=16, color=DARK, bold=(k == "quote"))
            cy += Inches(0.6)
    add_text(slide, MARGIN, Inches(6.9), Inches(12.2), Inches(0.4),
             "GitHub 开源 · 双击即跑 · 欢迎试用", size=14, color=GREEN, bold=True)


def render_content(prs, page, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide)
    title_bar(slide, page["title"], page_no, total)
    layout = page["dir"].get("layout", "")
    visual = page["dir"].get("visual", "")
    if layout == "two-col":
        lw = Inches(6.35)
        rx = Inches(7.05)
        rw = Inches(5.75)
        skip = False
    elif visual:
        lw = Inches(12.2)
        rx = MARGIN
        rw = Inches(12.2)
        skip = True  # 全宽页：示意图放底部
    else:
        lw = Inches(12.2)
        rx = MARGIN
        rw = Inches(12.2)
        skip = True
    consumes = visual.partition(":")[0] in VISUAL_CONSUMES_LI
    cy = render_items(slide, page, MARGIN, BODY_TOP, lw, BODY_H,
                      skip_li=consumes, size=15 if layout == "two-col" else 16)
    if visual:
        if layout == "two-col":
            draw_visual(slide, page, rx, BODY_TOP, rw, BODY_H)
        else:
            vh = Inches(3.0)
            vy = Inches(4.05)
            if consumes:
                vy = BODY_TOP + Inches(0.1)
                vh = Inches(4.3)
            draw_visual(slide, page, rx, vy, rw, vh)
    return slide


def build(md_path, out_path):
    md = Path(md_path).read_text(encoding="utf-8")
    pages = parse_pages(md)
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    total = len(pages)
    for i, page in enumerate(pages):
        ptype = page["dir"].get("type", "")
        if ptype == "cover":
            render_cover(prs, page, total)
        elif ptype == "thanks":
            render_thanks(prs, page, total)
        else:
            render_content(prs, page, i + 1, total)
    prs.save(out_path)
    print(f"OK: {out_path} ({total} slides)")


if __name__ == "__main__":
    md = sys.argv[1] if len(sys.argv) > 1 else "presentation.md"
    out = sys.argv[2] if len(sys.argv) > 2 else "Citrus_QA_Agent_汇报.pptx"
    build(md, out)
